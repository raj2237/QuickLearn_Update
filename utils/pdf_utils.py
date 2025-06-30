from fpdf import FPDF
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor
from reportlab.lib.units import inch
import fitz
from langchain_groq import ChatGroq
from config import Config
import re
from pptx import Presentation
import os
import aiofiles

async def extract_text_from_pdf(pdf_file):
    reader = fitz.open(pdf_file)
    text = " ".join(page.get_text("text") for page in reader if page.get_text("text"))
    reader.close()
    return text

async def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
    return " ".join(text)

async def extract_questions_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        questions = []
        current_question = ""
        
        for page in doc:
            text = page.get_text("text")
            if not text.strip():
                text = page.get_text("blocks")
                text = "\n".join(block[4] for block in text if len(block) > 4)
            
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            
            for line in lines:
                is_new_question = False
                if re.match(r'^\d+[.\)]\s', line):
                    is_new_question = True
                elif '?' in line and len(line) > 10:
                    is_new_question = True
                
                if is_new_question:
                    if current_question:
                        questions.append(current_question.strip())
                    current_question = line
                elif current_question:
                    current_question += " " + line
            
            if current_question:
                questions.append(current_question.strip())
                current_question = ""
        
        filtered_questions = []
        for q in questions:
            if '?' in q and len(q) > 10:
                cleaned_q = " ".join(q.split())
                filtered_questions.append(cleaned_q)
        
        if not filtered_questions:
            all_text = ""
            for page in doc:
                all_text += page.get_text("text") + "\n"
            potential_questions = [q.strip() for q in all_text.split('?') if q.strip()]
            for pq in potential_questions:
                if len(pq) > 10:
                    cleaned_q = f"{pq}?".strip()
                    filtered_questions.append(cleaned_q)
        
        doc.close()
        return filtered_questions if filtered_questions else ["No valid questions found in PDF"]
    except Exception as e:
        return [f"Error extracting questions: {str(e)}"]

async def generate_questions(extracted_questions, num_questions):
    llm = ChatGroq(
        model="llama-3.1-8b-instant",
        temperature=0.7,
        groq_api_key=Config.GROQ_API_KEY
    )
    
    prompt = (
        "Based on the following sample questions, generate similar academic questions that are clear, meaningful, "
        "and properly formatted. Each question should end with a question mark and be suitable for an exam paper:\n\n"
        f"Sample questions:\n{chr(10).join(extracted_questions[:5])}\n\n"
        f"Generate {num_questions} new questions in a similar style and complexity level."
    )
    
    response = await llm.ainvoke(prompt)
    generated_text = response.content
    
    new_questions = [q.strip() for q in generated_text.split('\n') if q.strip()]
    valid_questions = []
    for q in new_questions:
        if q and q[0].isdigit() and '?' in q and len(q) > 15:
            valid_questions.append(q)
        elif q and '?' in q and len(q) > 15:
            valid_questions.append(f"{len(valid_questions) + 1}. {q}")
    
    return valid_questions[:num_questions]

async def create_question_paper(questions, filename, set_number):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    pdf.set_fill_color(240, 248, 255)
    pdf.rect(0, 0, pdf.w, pdf.h, 'F')
    
    pdf.set_font("Arial", "B", 16)
    pdf.set_text_color(0, 0, 0)
    pdf.cell(0, 10, f"Set Number - {set_number}", ln=True, align='C')
    pdf.ln(10)
    
    pdf.set_font("Arial", size=12)
    for idx, question in enumerate(questions, 1):
        q_text = question.split('.', 1)[1].strip() if '.' in question[:3] else question
        safe_question = f"{idx}. {q_text}".encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, safe_question)
        pdf.ln(5)
    
    pdf.set_y(-20)
    pdf.set_font("Arial", "I", 8)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 10, "Powered by QuickLearn AI", ln=True, align='C')
    
    async with aiofiles.open(filename, 'wb') as f:
        await f.write(pdf.output(dest='S').encode('latin-1'))

async def create_question_bank_pdf(text, subject):
    question_bank_dir = os.path.join(os.getcwd(), "generated_papers")
    os.makedirs(question_bank_dir, exist_ok=True)
    
    filename = os.path.join(question_bank_dir, f"{subject.replace(' ', '_')}_Questions.pdf")
    
    try:
        doc = SimpleDocTemplate(filename, pagesize=letter,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=72)
        
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=28,
            spaceAfter=20,
            alignment=1,
            textColor=HexColor('#FFD700'),
            leading=32
        )
        
        question_style = ParagraphStyle(
            'CustomQuestion',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=12,
            leading=18,
            spaceBefore=20,
            spaceAfter=20,
            firstLineIndent=0,
            leftIndent=20,
            textColor=HexColor('#00FF9D')
        )
        
        elements = []
        elements.append(Paragraph(f"{subject} Practice Questions", title_style))
        elements.append(Spacer(1, 0.2*inch))
        
        questions = [q.strip() for q in text.split('\n\n') if q.strip()]
        for q in questions:
            q = q.replace('&', '&').replace('<', '<').replace('>', '>')
            p = Paragraph(q, question_style)
            elements.append(p)
        
        def add_watermark_and_background(canvas, doc):
            canvas.saveState()
            canvas.setFillColorRGB(0, 0, 0)
            canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=1)
            canvas.setFont("Helvetica-Bold", 60)
            canvas.setFillColorRGB(0.9, 0.9, 0.9, alpha=0.5)
            canvas.translate(doc.pagesize[0]/2, doc.pagesize[1]/2)
            canvas.rotate(45)
            canvas.drawCentredString(0, 0, "QuickLearn AI")
            canvas.restoreState()
        
        doc.build(elements, onFirstPage=add_watermark_and_background, onLaterPages=add_watermark_and_background)
        
        return filename
    
    except Exception as e:
        raise Exception(f"Error building PDF: {str(e)}")